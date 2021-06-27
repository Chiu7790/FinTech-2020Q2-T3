import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat, RGBColor
from PIL import Image
import os
import psutil
#from app_ppt import (cover, ppt_insert_summarization, ppt_insert_images)

# Ref for slide types: 
# 0 ->  title and subtitle
# 1 ->  title and content
# 2 ->  section header
# 3 ->  two content
# 4 ->  Comparison
# 5 ->  Title only 
# 6 ->  Blank
# 7 ->  Content with caption
# 8 ->  Pic with caption

# 首頁PPT
def cover(ppt_file, insert_title, insert_author, image, Layout=0, Placeholder=1, start_ppt=False):
    if os.path.exists(ppt_file):
        os.remove(ppt_file)
    prs=Presentation()

    title_slide_layout = prs.slide_layouts[Layout] #建立簡報檔頁面物件
    #使用簡報物件中的方法將上一行建立的頁面物件放進簡報
    slide = prs.slides.add_slide(title_slide_layout)
    #設定頁面的標題 
    title = slide.shapes.title
    title.text = insert_title
    #設定頁面的副標題
    subtitle = slide.placeholders[Placeholder] #設定副標題物件，副標題通常為第2個佔位圖
    subtitle.text = insert_author       

    # 首頁圖片
    height = Inches(2)
    left = Inches(4)
    top = Inches(5)
    pic = slide.shapes.add_picture(image, left, top, height=height)

    #將簡報物件存檔
    prs.save(ppt_file)
    if start_ppt:
        os.startfile(ppt_file)

# 目錄頁
def contents(ppt_file, insert_title, paragraph, Layout=1, Placeholder=1, start_ppt=False):
    if os.path.exists(ppt_file):
        prs = Presentation(ppt_file)
    else:
        prs=Presentation()

    title_slide_layout = prs.slide_layouts[Layout]  # 建立簡報檔頁面物件
    #使用簡報物件中的方法將上一行建立的頁面物件放進簡報
    slide = prs.slides.add_slide(title_slide_layout)

    #設定頁面的標題
    title = slide.shapes.title
    title.text = insert_title[0]

    #設定頁面的內容
    contents = slide.placeholders[Placeholder]  # 設定副標題物件，副標題通常為第2個佔位圖
    # creating textFrames
    tf = contents.text_frame

    # adding Paragraphs
    for i in range(len(paragraph)):
        p = tf.add_paragraph()
        p.text = paragraph[i]
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(30, 144, 255)

    #將簡報物件存檔
    prs.save(ppt_file)
    if start_ppt:
        os.startfile(ppt_file)

# 中間分隔頁
def divider(ppt_file, paragraph, insert_sub, divide_page, Layout=2, start_ppt=False):
    if os.path.exists(ppt_file):
        prs = Presentation(ppt_file)
    else:
        prs = Presentation()

    title_slide_layout = prs.slide_layouts[Layout]  # 建立簡報檔頁面物件
    #使用簡報物件中的方法將上一行建立的頁面物件放進簡報
    slide = prs.slides.add_slide(title_slide_layout)
    #設定頁面的標題
    title = slide.shapes.title
    title.text = paragraph[divide_page-1]
    #設定頁面的內容
    subtitle = slide.placeholders[1]  # 設定副標題物件，副標題通常為第2個佔位圖
    subtitle.text = insert_sub[0]

    #將簡報物件存檔
    prs.save(ppt_file)
    if start_ppt:
        os.startfile(ppt_file)

# words：Hot Words、Related Words
def words(ppt_file, insert_title, paragraph, font_size, Layout=3, start_ppt=False):
    if os.path.exists(ppt_file):
        prs = Presentation(ppt_file)
    else:
        prs = Presentation()

    title_slide_layout = prs.slide_layouts[Layout]  # 建立簡報檔頁面物件
    #使用簡報物件中的方法將上一行建立的頁面物件放進簡報
    slide = prs.slides.add_slide(title_slide_layout)

    #設定頁面的標題
    title = slide.shapes.title
    title.text = insert_title[0]

    #設定頁面的內容
    contents = slide.placeholders[1]  # 設定左側物件，副標題通常為第2個佔位圖
    # creating textFrames
    tf = contents.text_frame
    tf.text = ""
    # adding Paragraphs
    for i in range(int(len(paragraph)/2)):
        p = tf.add_paragraph()
        p.text = paragraph[i]
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(30, 144, 255)
    
    #設定第三張頁面的內容
    contents = slide.placeholders[2]  # 設定右側物件，副標題通常為第3個佔位圖
    # creating textFrames
    tf = contents.text_frame
    tf.text = ""
    # adding Paragraphs
    for i in range(int(len(paragraph)/2), int(len(paragraph))):
        p = tf.add_paragraph()
        p.text = paragraph[i]
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(30, 144, 255)

    #將簡報物件存檔
    prs.save(ppt_file)
    if start_ppt:
        os.startfile(ppt_file)

# Only Picture：Word Clouds、Knowledge Graph Co-occurence-matrix、Stock
def only_pic(ppt_file, pic_title, image, left_inch, top_inch, size, Layout=1, start_ppt=False):
    if os.path.exists(ppt_file):
        prs = Presentation(ppt_file)
    else:
        prs=Presentation()

    title_slide_layout = prs.slide_layouts[Layout]  # 建立簡報檔頁面物件
    #使用簡報物件中的方法將上一行建立的頁面物件放進簡報
    slide = prs.slides.add_slide(title_slide_layout)

    #設定頁面的標題
    title = slide.shapes.title  # 設定左側物件，副標題通常為第2個佔位圖
    title.text = pic_title

    # # convert pixels to inches
    # def px_to_inches(image):
    #     im = Image.open(image)
    #     width = im.width / im.info['dpi'][0]
    #     height = im.height / im.info['dpi'][1]
    #     return (width, height)
    # img = px_to_inches(image)

    # # show the figure
    # left = int((prs.slide_width - img[0]) / 2)
    # top= int((prs.slide_height - img[1]) / 2)
    height = Inches(size)
    left = Inches(left_inch)
    top = Inches(top_inch)
    pic = slide.shapes.add_picture(image, left, top, height=height)
  
    #將簡報物件存檔
    prs.save(ppt_file)
    if start_ppt:
        os.startfile(ppt_file)

# Only text
def only_text(ppt_file, insert_title, summary, Layout=1, Placeholder=1, start_ppt=False):
    if os.path.exists(ppt_file):
        prs = Presentation(ppt_file)
    else:
        prs=Presentation()

    title_slide_layout = prs.slide_layouts[Layout]  # 建立簡報檔頁面物件
    #使用簡報物件中的方法將上一行建立的頁面物件放進簡報
    slide = prs.slides.add_slide(title_slide_layout)

    #設定頁面的標題
    title = slide.shapes.title
    title.text = insert_title

    #設定頁面的內容
    contents = slide.placeholders[Placeholder]  # 設定副標題物件，副標題通常為第2個佔位圖
    # creating textFrames
    tf = contents.text_frame

    p = tf.add_paragraph()
    p.text = summary[0]
    p.font.size = Pt(16)

    #將簡報物件存檔
    prs.save(ppt_file)
    if start_ppt:
        os.startfile(ppt_file)

# CAPM
def text_picture(ppt_file, text_pic_title, paragraph, image, left_inch, top_inch, size, Layout=3, start_ppt=False):
    if os.path.exists(ppt_file):
        prs = Presentation(ppt_file)
    else:
        prs = Presentation()

    title_slide_layout = prs.slide_layouts[Layout]  # 建立簡報檔頁面物件
    #使用簡報物件中的方法將上一行建立的頁面物件放進簡報
    slide = prs.slides.add_slide(title_slide_layout)

    #設定頁面的標題
    title = slide.shapes.title
    title.text = text_pic_title

    #設定頁面的內容
    contents = slide.placeholders[1]  # 設定左側物件，副標題通常為第2個佔位圖
    # creating textFrames
    tf = contents.text_frame
    tf.text = ""
    # adding Paragraphs
    for i in range(len(paragraph)):
        p = tf.add_paragraph()
        p.text = paragraph[i]
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(30, 144, 255)

    #設定第三張頁面的內容-pic
    height = Inches(size)
    left = Inches(left_inch)
    top = Inches(top_inch)
    pic = slide.shapes.add_picture(image, left, top, height=height)

    #將簡報物件存檔
    prs.save(ppt_file)
    if start_ppt:
        os.startfile(ppt_file)

# read file for readlines without \n
def read_file(path):
    with open(path, "r", encoding='utf-8') as f:
        lines = f.read().splitlines()
        return lines #list

# main app
def app(ppt_file='find_trend.pptx'):
    for proc in psutil.process_iter():
        if proc.name() == 'POWERPNT.EXE':
            proc.kill()

    # Home
    with open('report/cover/title.txt', mode='r', encoding='utf-8') as f:
        insert_title = f.read()
    with open('report/cover/subtitle.txt', mode='r', encoding='utf-8') as f:
        insert_author = f.read()
    image_path = 'report/cover/img/trending.png'

    cover(ppt_file=ppt_file, insert_title=insert_title,
          insert_author=insert_author, image = image_path, start_ppt=False)

    # Content
    content_title = read_file('report/contents/title.txt')
    paragraph = read_file('report/contents/contents.txt')
    contents(ppt_file=ppt_file, insert_title = content_title, paragraph=paragraph, start_ppt=False)

    # Divider1-Hot Words
    insert_sub = read_file('report/hot_word/subtitle.txt')
    # 透過Google trend選取近期討論度較高的字詞後，讓使用者可自行決定使否要查詢熱門字詞的相關新聞及主題摘要，抑或是自己主動搜尋。
    divider(ppt_file, paragraph,
            insert_sub=insert_sub, divide_page=1, start_ppt=False)

    # Hot Words
    hotwords_title = read_file('report/hot_word/title.txt')
    hot_words = read_file('report/hot_word/trendword.txt')
    words(ppt_file=ppt_file, insert_title=hotwords_title,
          paragraph=hot_words, font_size=20, start_ppt=False)

    # Divider2-Keyword
    insert_sub = read_file('report/keyword/subtitle.txt')
    # 針對關鍵字爬取特定詞的新聞，並透過摘要自動生成模組，輸出可以調整文章量多寡的摘要。
    divider(ppt_file, paragraph,
            insert_sub=insert_sub, divide_page=2, start_ppt=False)
        
    # wordcloud
    keyword = read_file('report/keyword/keyword.txt')
    pic_title = "Word Cloud of【" + keyword[0] + "】News"
    image_path = 'report\keyword\img\wordcloud.png'
    only_pic(ppt_file, pic_title, image=image_path, left_inch=2.5, top_inch=2, size=5, start_ppt=False)

    # Summary
    insert_title = "Summary of 【" + keyword[0] + "】Keyword News"
    summary = read_file('report/keyword/keyword_summary.txt')
    only_text(ppt_file, insert_title, summary = summary, start_ppt=False)

    # Divider3-Relatedword
    insert_sub = read_file('report/related_word/subtitle.txt')
    # 針對關鍵字爬取特定詞的新聞，再經由知識圖譜模型取得相關的關聯詞。
    divider(ppt_file, paragraph,
            insert_sub=insert_sub, divide_page=3, start_ppt=False)

    # Related Words
    related_title = read_file('report/Related_word/title.txt')
    related_words = read_file('report/Related_word/relatedword.txt')
    words(ppt_file=ppt_file, insert_title=related_title,
          paragraph=related_words, font_size=18, start_ppt=False)

    # Knowledge Graph
    pic_title = "Knowledge Graph of 【" + keyword[0] + "】News"
    image_path = 'report\Related_word\img\knowledge_graph.png'
    only_pic(ppt_file, pic_title, image=image_path,
             left_inch=1.5, top_inch=1, size=7, start_ppt=False)

    # Co-occurence matrix
    pic_title = "Co-occurence matrix of 【" + keyword[0] + "】News"
    image_path = 'report\Related_word\img\co_occurence_matrix.png'
    only_pic(ppt_file, pic_title, image=image_path,
             left_inch=2.5, top_inch=2, size=5, start_ppt=False)
        
    # Divider4-Stock
    insert_sub = read_file('report/stock/subtitle.txt')
    # 畫出各檔股價的股價趨勢圖及利用美國1年期公債殖利率作為無風險利率的報酬來計算CAPM的預期一年報酬率。
    divider(ppt_file, paragraph,
            insert_sub=insert_sub, divide_page=4, start_ppt=False)

    # stock trend chart
    target_stock = read_file('report/stock/target_stock.txt')
    pic_title = "【" + target_stock[0] + "】Price Trend Chart"
    image_path = 'report/stock/img/stock_ticker.png'
    only_pic(ppt_file, pic_title, image=image_path,
             left_inch=1, top_inch=2, size=5, start_ppt=False) #位置要調整
    
    # Price Trend Chat
    pic_title = "【" + target_stock[0] + "】Price Prediction Trend Chat"
    image_path = 'report/stock/img/fbprophet_ticker.png'
    only_pic(ppt_file, pic_title, image=image_path,
             left_inch=1, top_inch=2, size=5, start_ppt=False)  # 位置要調整

    # CAPM
    text_pic_title = "【" + target_stock[0] + "】CAPM Value"
    image_path = 'report/stock/img/capm.png'
    capm_ls = read_file('report/stock/stock_ls.txt')
    text_picture(ppt_file, text_pic_title, paragraph=capm_ls, image=image_path,
                 left_inch = 3.5, top_inch = 2, size = 4, start_ppt=True)
    
# if st.button("生成簡報檔"):
#     app()
